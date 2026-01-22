import logging
from pathlib import Path
from typing import List
from uuid import UUID
from pptx import Presentation

logger = logging.getLogger(__name__)

class PptxParser:
    """
    Handles extraction of text from PPTX files.
    """

    async def parse(self, session_id: UUID, file_path: Path, output_dir: Path) -> List[str]:
        """
        Parses a PPTX file and returns a list of text chunks (one per slide).
        """
        if not file_path.exists():
            raise FileNotFoundError(f"PPTX not found: {file_path}")

        logger.info(f"[{session_id}] Parsing PPTX: {file_path.name}")
        
        try:
            prs = Presentation(file_path)
            chunks = []
            
            for i, slide in enumerate(prs.slides):
                slide_content = [f"--- Slide {i+1} ---"]
                
                # Extract text from shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text = shape.text.strip()
                        if text:
                            slide_content.append(text)
                            
                # Extract speaker notes (crucial for context)
                if slide.has_notes_slide:
                    notes = slide.notes_slide.notes_text_frame.text.strip()
                    if notes:
                        slide_content.append(f"[NOTES]: {notes}")
                        
                chunks.append("\n".join(slide_content))
                
            return chunks

        except Exception as e:
            logger.error(f"[{session_id}] PPTX Parse error: {e}")
            raise RuntimeError(f"Failed to parse PPTX: {e}")
